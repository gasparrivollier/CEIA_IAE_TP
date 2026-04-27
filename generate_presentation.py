from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Crear presentación
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Colores personalizados
COLOR_DARK_BLUE = RGBColor(30, 39, 97)      # #1E2761
COLOR_LIGHT_BLUE = RGBColor(202, 220, 252)  # #CADCFC
COLOR_ACCENT = RGBColor(255, 255, 255)      # Blanco
COLOR_SECONDARY = RGBColor(28, 114, 147)    # #1C7293
COLOR_TEXT = RGBColor(50, 50, 50)
COLOR_LIGHT_BG = RGBColor(245, 245, 245)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BLUE
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = COLOR_LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items, has_bg=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    if has_bg:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_LIGHT_BG
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_BLUE
    
    # Línea decorativa
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = COLOR_SECONDARY
    line.line.width = Pt(3)
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

def add_two_column_slide(prs, title, left_items, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_BG
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_BLUE
    
    # Línea decorativa
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = COLOR_SECONDARY
    line.line.width = Pt(3)
    
    # Columna izquierda
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.2))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)
    
    # Columna derecha
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(5.2))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)

# ========== DIAPOSITIVAS ==========

# 1. Portada
add_title_slide(prs, 
    "Optimización y Despliegue de IA\nen Entornos Restringidos",
    "Clasificación de Lenguaje de Señas Americano (ASL)\ncon Compresión de Modelos Deep Learning")

# 2. Índice
add_content_slide(prs, "Contenido", [
    "✓ Fase 1: Definición del Problema y Restricciones",
    "✓ Fase 2: Arquitectura Baseline y Entrenamiento",
    "✓ Fase 3: Técnicas de Optimización del Modelo",
    "✓ Fase 4: Despliegue y Benchmarking",
    "✓ Resultados Comparativos",
    "✓ Conclusiones"
])

# 3. Fase 1 - Problema
add_content_slide(prs, "Fase 1: Definición del Problema", [
    "📌 Problema: Clasificación de Lenguaje de Señas (ASL)",
    "   • Dataset: ASL Alphabet (Kaggle - grassknoted/asl-alphabet)",
    "   • Clases: 29 señas del alfabeto americano",
    "   • Tipo de tarea: Clasificación de imágenes multiclase",
    "",
    "🎯 Aplicación: Traducción automática de lenguaje de señas",
    "   en tiempo real para dispositivos con restricciones de hardware"
])

# 4. Fase 1 - Restricciones
add_content_slide(prs, "Fase 1: Restricciones de Hardware", [
    "💾 Restricción de Memoria: ≤ 100 MB en disco",
    "   • Modelo baseline FP32: 9.5 MB",
    "",
    "⚡ Restricción de Latencia: < 50 ms por inferencia",
    "   • En CPU (entorno restringido)",
    "",
    "🎯 Restricción de Precisión: ≥ 77% accuracy (drop ≤ 2%)",
    "   • Mantener funcionalidad del modelo",
    "",
    "📊 Métrica de Éxito: Reducir tamaño ≥50% sin perder precisión"
])

# 5. Fase 2 - Arquitectura
add_content_slide(prs, "Fase 2: Arquitectura CNN Baseline", [
    "🏗️ Diseño: CNN personalizada (4 bloques convolucionales)",
    "",
    "📐 Capas:",
    "   • Conv1: 32 filtros 3×3 + BatchNorm + ReLU + MaxPool",
    "   • Conv2: 64 filtros 3×3 + BatchNorm + ReLU + MaxPool", 
    "   • Conv3: 128 filtros 3×3 + BatchNorm + ReLU + MaxPool",
    "   • Conv4: 256 filtros 3×3 + BatchNorm + ReLU + MaxPool",
    "   • FC: 512 → 256 → 29 clases",
    "",
    "🔧 Regularización: Dropout, BatchNormalization"
])

# 6. Métricas Baseline
add_two_column_slide(prs, "Fase 2: Métricas Baseline (FP32)",
    [
        "📊 Complejidad Teórica:",
        "• Parámetros: 2.49 M",
        "• MACs: 63.25 M",
        "• Tamaño: 9.5 MB",
        "",
        "⚙️ Entrenamiento:",
        "• 25 épocas",
        "• Optimizador: Adam",
        "• Loss: CrossEntropyLoss"
    ],
    [
        "✅ Performance (CPU):",
        "• Val Accuracy: 79.32%",
        "• Test Accuracy: 79.09%",
        "• Val Loss: 0.9008",
        "",
        "⏱️ Inferencia:",
        "• Latencia: 4.98 ms",
        "• FPS: ~201",
        "• Peak Memory: ~350 MB"
    ])

# 7. Fase 3 - Técnicas de Optimización
add_content_slide(prs, "Fase 3: Técnicas de Optimización", [
    "🔧 Cuantización: Conversión de FP32 → INT8/INT4",
    "   • Dinámica (dynamic quantization)",
    "   • Estática (static quantization)",
    "   • Quantization-Aware Training (QAT)",
    "",
    "✂️ Poda (Pruning): Eliminación de pesos no significativos",
    "   • Ratio de poda: 40%",
    "",
    "📦 Factorización: Descomposición de convoluciones",
    "   • Separación en convoluciones depthwise + pointwise",
    "",
    "🎓 Destilación: Transfer de conocimiento (Teacher-Student)"
])

# 8. Resultados Cuantización
add_two_column_slide(prs, "Fase 3: Resultados - Cuantización",
    [
        "INT8 Dinámica:",
        "• Accuracy: 79.15%",
        "• Tamaño: 3.51 MB (-63%)",
        "• Latencia: 4.85 ms",
        "",
        "INT8 Estática:",
        "• Accuracy: 79.22%",
        "• Tamaño: 2.38 MB (-75%)",
        "• Latencia: 3.92 ms"
    ],
    [
        "✅ Beneficios:",
        "• Reducción: 63-75%",
        "• Sin pérdida de accuracy",
        "• Compatible CPU",
        "",
        "⚡ Mejora de velocidad:",
        "• 20-27% más rápido",
        "• Menor consumo RAM"
    ])

# 9. Resultados Poda
add_content_slide(prs, "Fase 3: Resultados - Pruning (40%)", [
    "✂️ Modelo Podado (FP32):",
    "   • Parámetros: 1.47 M (-41%)",
    "   • Tamaño: 5.6 MB (-41%)",
    "   • Accuracy: 78.95% (-0.14%)",
    "   • Latencia: 3.52 ms (-29%)",
    "",
    "✂️ Podado + INT8 Estático (Óptimo):",
    "   • Tamaño: 1.42 MB (-85%)",
    "   • Accuracy: 79.11% (preservado)",
    "   • Latencia: 2.15 ms (-57%)",
    "   • Cumple todas las restricciones ✓"
])

# 10. Resultados Distilación
add_content_slide(prs, "Fase 3: Resultados - Knowledge Distillation", [
    "🎓 Student Model (MLP comprimido):",
    "   • Parámetros: 315K (-87%)",
    "   • Tamaño: 1.21 MB (-87%)",
    "   • Accuracy: 77.45% (-1.64%)",
    "",
    "🎓 Student + INT8 Estático:",
    "   • Tamaño: 0.31 MB (-97%)",
    "   • Accuracy: 77.38% (dentro del límite)",
    "   • Latencia: 1.05 ms (-79%)",
    "   • Máxima compresión alcanzada"
])

# 11. Tabla Comparativa
add_content_slide(prs, "Fase 4: Comparativa de Modelos", [
    "┌─────────────────────────────────────────────────────┐",
    "│ Modelo         │ Tamaño  │ Accuracy │ Latencia │    │",
    "├─────────────────────────────────────────────────────┤",
    "│ Baseline FP32  │ 9.5 MB  │ 79.09%   │ 4.98 ms  │ 📊 │",
    "│ INT8 Dinámico  │ 3.51 MB │ 79.15%   │ 4.85 ms  │    │",
    "│ INT8 Estático  │ 2.38 MB │ 79.22%   │ 3.92 ms  │    │",
    "│ Podado 40%     │ 5.6 MB  │ 78.95%   │ 3.52 ms  │    │",
    "│ Pod. + INT8    │ 1.42 MB │ 79.11%   │ 2.15 ms  │ ⭐ │",
    "│ Student INT8   │ 0.31 MB │ 77.38%   │ 1.05 ms  │ 🚀 │",
    "└─────────────────────────────────────────────────────┘"
])

# 12. Análisis de Resultados
add_content_slide(prs, "Fase 4: Análisis de Resultados", [
    "✅ Restricciones Cumplidas:",
    "   ✓ Tamaño ≤ 100 MB (alcanzado 0.31-1.42 MB)",
    "   ✓ Latencia < 50 ms (alcanzado 1.05-2.15 ms)",
    "   ✓ Drop de precision ≤ 2% (máximo 1.64%)",
    "",
    "📈 Logros Principales:",
    "   • Reducción de 85-97% en tamaño",
    "   • Mejora de 20-79% en velocidad",
    "   • Preservación de precisión del modelo"
])

# 13. Conclusiones
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_DARK_BLUE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Conclusiones"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT
p.alignment = PP_ALIGN.CENTER

content_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(3.8))
text_frame = content_box.text_frame
text_frame.word_wrap = True

items = [
    "✓ Compresión exitosa: 85-97% reducción de tamaño",
    "✓ Modelos optimizados viables para dispositivos embebidos",
    "✓ Combinación Pruning + Cuantización es óptima",
    "✓ Knowledge Distillation logra máxima compresión",
    "✓ Ciclo completo de optimización completado exitosamente"
]

for i, item in enumerate(items):
    if i > 0:
        text_frame.add_paragraph()
    p = text_frame.paragraphs[i]
    p.text = item
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_LIGHT_BLUE
    p.space_before = Pt(10)

# Guardar
prs.save('TP_Integrador_Optimizacion_IA.pptx')
print("✓ Presentación creada exitosamente: TP_Integrador_Optimizacion_IA.pptx")
